import os
import fitz
import docx
from pptx import Presentation
from openai import OpenAI
from dotenv import load_dotenv
import json
import re

# Load environment variables
load_dotenv()
api_key = os.getenv("OPENROUTER_API_KEY")

# Initialize OpenAI client for OpenRouter
client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=api_key,
)

def extract_text_from_file(filepath: str) -> str:
    ext = os.path.splitext(filepath)[-1].lower()
    if ext == '.txt':
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    elif ext == '.pdf':
        doc = fitz.open(filepath)
        return "\n".join(page.get_text() for page in doc)
    elif ext == '.docx':
        doc = docx.Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    elif ext == '.pptx':
        prs = Presentation(filepath)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    else:
        raise ValueError("Unsupported file format")

def extract_json_string(response_str):
    json_match = re.search(r'{\s*".*?"\s*:\s*\[.*?\]\s*}|{\s*".*?"\s*:\s*".*?"\s*}', response_str, re.DOTALL)
    if json_match:
        return json_match.group()
    return None

def generate_quiz_and_answers(text: str) -> dict:
    prompt = """
    You are an expert educational assistant tasked with creating a quiz based on the content of a user-uploaded file, typically a presentation (e.g., PowerPoint, PDF) provided by university professors.
    Your goal is to generate a short quiz to test understanding of the key concepts, facts, and ideas in the presentation.
    Follow these guidelines:
    1. Analyze the Content:
       - Carefully interpret the provided content.
       - Identify main topics, key concepts, definitions, theories, and examples.
       - If visuals are mentioned, infer their meaning if possible.
    2. Generate Quiz Questions:
       - Create 5 to 25 questions depending on the content’s depth.
       - Use a mix of formats:
         - Multiple-choice (include exactly 4 options)
         - True/False
         - Short-answer
         - Fill-in-the-blank
       - Keep questions clear, relevant, and educational.
       - Ensure each question is a single, concise string.
    3. Format:
       - Number each question clearly: "Question 1", "Question 2", etc.
       - Provide only the list of questions.
       - Do NOT include the answers in the question text.
    4. Store Correct Answers Internally:
       - Prepare the correct answers in order.
       - Each answer should be a single, concise string.
       - Only show them when the user requests them (e.g., "Show answers").
    5. Edge Cases:
       - If the content is too short, generate as many meaningful questions as possible (minimum 3).
       - If content is unclear, make a reasonable assumption.
    6. Style:
       - Use a professional academic tone.
       - Make the quiz engaging and challenging but fair.
    **Important: Only return a JSON object with the following exact structure:**
    ```json
    {
      "quiz": ["Question 1: ...", "Question 2: ...", "..."],
      "answers": ["Answer 1: ...", "Answer 2: ...", "..."]
    }
    ```
    Ensure each question and answer is a single string, even for complex answers.
    """
    try:
        response = client.chat.completions.create(
            model="meta-llama/llama-4-maverick:free",
            messages=[{"role": "user", "content": prompt + "\n\n" + text}],
            temperature=0.7
        )
        raw_content = response.choices[0].message.content.strip()
        json_str = extract_json_string(raw_content)
        if not json_str:
            raise ValueError("No valid JSON object found in the response.")
        data = json.loads(json_str)
        if not isinstance(data, dict) or "quiz" not in data or "answers" not in data:
            raise ValueError("Response does not contain 'quiz' and 'answers' keys.")
        quiz = [str(q).strip() for q in data["quiz"]]
        answers = [str(a).strip() for a in data["answers"]]
        if len(quiz) != len(answers):
            raise ValueError(f"Mismatch between number of questions ({len(quiz)}) and answers ({len(answers)}).")
        return {"quiz": quiz, "answers": answers}
    except Exception as e:
        print(f"Error generating quiz: {e}")
        return None

def generate_notes(text: str) -> str:
    prompt = """
    You are an expert educational assistant tasked with creating concise study notes based on the content of a user-uploaded file, typically a presentation (e.g., PowerPoint, PDF) provided by university professors.
    Your goal is to generate a short summary with bullet points highlighting the key concepts, facts, definitions, theories, and examples in the content.
    Follow these guidelines:
    1. Analyze the Content:
       - Carefully interpret the provided content.
       - Identify main topics, key concepts, definitions, theories, and examples.
       - If visuals are mentioned, infer their meaning if possible.
    2. Generate Notes:
       - Provide a brief introductory summary (1-2 sentences) of the content.
       - List 5-10 bullet points covering the most important points.
       - Keep each bullet point concise and focused (1-2 sentences).
       - Use clear, academic language suitable for study purposes.
    3. Format:
       - Return a JSON object with a single key "notes" containing a list of strings.
       - The first string is the summary paragraph.
       - Subsequent strings are bullet points starting with "- ".
    4. Edge Cases:
       - If the content is too short, generate as many meaningful bullet points as possible (minimum 3).
       - If content is unclear, make reasonable assumptions.
    5. Style:
       - Use a professional academic tone.
       - Ensure the notes are clear, organized, and useful for studying.
    **Important: Only return a JSON object with the following exact structure:**
    ```json
    {
      "notes": ["Summary: ...", "- Bullet point 1", "- Bullet point 2", "..."]
    }
    ```
    Ensure each note (summary and bullet points) is a single string.
    """
    try:
        response = client.chat.completions.create(
            model="meta-llama/llama-4-maverick:free",
            messages=[{"role": "user", "content": prompt + "\n\n" + text}],
            temperature=0.7
        )
        raw_content = response.choices[0].message.content.strip()
        json_str = extract_json_string(raw_content)
        if not json_str:
            raise ValueError("No valid JSON object found in the response.")
        data = json.loads(json_str)
        if not isinstance(data, dict) or "notes" not in data:
            raise ValueError("Response does not contain 'notes' key.")
        notes = [str(n).strip() for n in data["notes"]]
        if not notes:
            raise ValueError("No notes generated.")
        return "\n\n".join(notes)
    except Exception as e:
        print(f"Error generating notes: {e}")
        return None

def generate_mnemonics(words: list) -> str:
    prompt = """
    You are an expert educational assistant tasked with creating mnemonic aids for a list of user-provided words to help with memorization.
    Your goal is to generate creative and effective mnemonics using techniques such as acronyms, associations, imagery, or rhymes.
    Follow these guidelines:
    1. Analyze the Words:
       - Consider the meaning, spelling, or context of each word if applicable.
       - Use the word itself to inspire the mnemonic.
    2. Generate Mnemonics:
       - Create one mnemonic per word.
       - Each mnemonic should be concise, memorable, and relevant to the word.
       - Use techniques like:
         - Acronyms (e.g., for 'CAT', "Cute Animal Tail").
         - Associations (e.g., for 'RIVER', "Imagine a flowing river carving a valley").
         - Imagery (e.g., for 'SUN', "Picture a bright sun warming your face").
         - Rhymes or phrases (e.g., for 'BOOK', "Look, a book to hook your mind").
       - Ensure the mnemonic is educational and appropriate.
    3. Format:
       - Return a JSON object with a single key "mnemonics" containing a list of strings.
       - Each string should be formatted as "Word: <word> - Mnemonic: <mnemonic>".
    4. Edge Cases:
       - If a word is complex or unclear, make a reasonable assumption about its meaning.
       - If the list is empty, return an empty list.
    5. Style:
       - Use a professional yet engaging tone.
       - Ensure mnemonics are clear and useful for studying.
    **Important: Only return a JSON object with the following exact structure:**
    ```json
    {
      "mnemonics": ["Word: word1 - Mnemonic: ...", "Word: word2 - Mnemonic: ...", "..."]
    }
    ```
    Ensure each mnemonic entry is a single string.
    """
    try:
        response = client.chat.completions.create(
            model="meta-llama/llama-4-maverick:free",
            messages=[{"role": "user", "content": prompt + "\n\nWords: " + ", ".join(words)}],
            temperature=0.7
        )
        raw_content = response.choices[0].message.content.strip()
        json_str = extract_json_string(raw_content)
        if not json_str:
            raise ValueError("No valid JSON object found in the response.")
        data = json.loads(json_str)
        if not isinstance(data, dict) or "mnemonics" not in data:
            raise ValueError("Response does not contain 'mnemonics' key.")
        mnemonics = [str(m).strip() for m in data["mnemonics"]]
        if not mnemonics and words:
            raise ValueError("No mnemonics generated.")
        return "\n\n".join(mnemonics)
    except Exception as e:
        print(f"Error generating mnemonics: {e}")
        return None

def generate_story(words: list) -> str:
    prompt = """
    You are a creative storytelling assistant tasked with generating an engaging story based on a list of user-provided words.
    Your goal is to create a coherent and imaginative narrative (300-500 words) that incorporates all the provided words in a meaningful way.
    Follow these guidelines:
    1. Analyze the Words:
       - Consider the meaning, context, or potential thematic connections of each word.
       - Use each word naturally within the story, ensuring it fits the narrative flow.
    2. Generate the Story:
       - Create a single, continuous narrative that uses all provided words at least once.
       - The story should have a clear beginning, middle, and end.
       - Make the story engaging, suitable for a general audience, and appropriate for educational use.
       - Use a creative and vivid writing style to captivate the reader.
    3. Format:
       - Return a JSON object with a single key "story" containing the story as a single string.
       - Do not include additional headings, labels, or formatting within the story text.
    4. Edge Cases:
       - If the word list is empty, return an empty string.
       - If a word is complex or unclear, make a reasonable assumption about its meaning.
       - If there are many words (e.g., >10), prioritize natural integration over forced usage.
    5. Style:
       - Use a professional yet creative tone.
       - Ensure the story is clear, engaging, and suitable for study or entertainment purposes.
    **Important: Only return a JSON object with the following exact structure:**
    ```json
    {
      "story": "..."
    }
    ```
    Ensure the story is a single string.
    """
    try:
        response = client.chat.completions.create(
            model="meta-llama/llama-4-maverick:free",
            messages=[{"role": "user", "content": prompt + "\n\nWords: " + ", ".join(words)}],
            temperature=0.7
        )
        raw_content = response.choices[0].message.content.strip()
        json_str = extract_json_string(raw_content)
        if json_str:
            data = json.loads(json_str)
            if not isinstance(data, dict) or "story" not in data:
                raise ValueError("Response does not contain 'story' key.")
            story = str(data["story"]).strip()
            if not story and words:
                raise ValueError("No story generated.")
            return story
        else:
            # Fallback: Treat raw_content as the story if JSON parsing fails
            if raw_content and len(raw_content) > 50:  # Ensure it's not just an error message
                return raw_content.strip()
            raise ValueError("No valid story content found in the response.")
    except Exception as e:
        print(f"Error generating story: {e}")
        return None

def main(choice: str, filepath: str = None, words: list = None):
    print("📂 Welcome to ChatGPT Study Tool!")
    try:
        if choice in ["quiz", "q"]:
            if not filepath or not os.path.isfile(filepath):
                print("❌ Error: File does not exist.")
                return None, None
            text = extract_text_from_file(filepath)
            qa = generate_quiz_and_answers(text)
            if qa is None:
                print("⚠️ Failed to generate quiz and answers.")
                return None, None
            quizzes = qa["quiz"]
            answers = qa["answers"]
            print("\n📝 Quiz:")
            for j, question in enumerate(quizzes, 1):
                print(f"{j}. {question}")
            return quizzes, answers
        elif choice in ["notes", "n"]:
            if not filepath or not os.path.isfile(filepath):
                print("❌ Error: File does not exist.")
                return None, None
            text = extract_text_from_file(filepath)
            notes = generate_notes(text)
            if notes is None:
                print("⚠️ Failed to generate notes.")
                return None, None
            print("\n📝 Study Notes:")
            print(notes)
            return notes, None
        elif choice in ["mnemonics", "m"]:
            if not words:
                print("❌ No words provided.")
                return None, None
            mnemonics = generate_mnemonics(words)
            if mnemonics is None:
                print("⚠️ Failed to generate mnemonics.")
                return None, None
            print("\n📝 Mnemonics:")
            print(mnemonics)
            return mnemonics, None
        elif choice in ["story", "s"]:
            if not words:
                print("❌ No words provided.")
                return None, None
            story = generate_story(words)
            if story is None:
                print("⚠️ Failed to generate story.")
                return None, None
            print("\n📝 Story:")
            print(story)
            return story, None
        else:
            print("❌ Invalid choice. Please enter 'quiz', 'notes', 'mnemonics', or 'story'.")
            return None, None
    except Exception as e:
        print(f"⚠️ An error occurred: {e}")
        return None, None

if __name__ == "__main__":
    choice = input("📋 Do you want to generate a Quiz, Notes, Mnemonics, or Story? (quiz/notes/mnemonics/story): ").strip().lower()
    if choice in ["quiz", "q", "notes", "n"]:
        file_path = input("📎 Enter the path to your file: ").strip()
        result, answers = main(choice, filepath=file_path)
    elif choice in ["mnemonics", "m", "story", "s"]:
        words_input = input("📚 Enter words to generate mnemonics or story for (comma-separated): ").strip()
        words = [word.strip() for word in words_input.split(",") if word.strip()]
        result, answers = main(choice, words=words)
    else:
        print("❌ Invalid choice. Please enter 'quiz', 'notes', 'mnemonics', or 'story'.")
        result, answers = None, None

    if result and isinstance(result, list) and answers:  # For quiz
        want_answers = input("\n🔍 Would you like to see the answers now? (yes/no): ").strip().lower()
        if want_answers in ["yes", "y"]:
            print("\n✅ Answers:")
            for i, a in enumerate(answers, 1):
                print(f"{i}. {a}")
        else:
            print("\n👍 You can come back to check the answers later.")
    elif result and isinstance(result, str):  # For notes, mnemonics, or story
        print("\n👍 Output generated successfully.")