import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from openai import OpenAI
from datetime import datetime

# ===== CONFIGURATION =====
OPENROUTER_API_KEY = "api-key"  # REPLACE THIS

# ===== INITIALIZATION =====
client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=OPENROUTER_API_KEY,
)

def save_to_file(content, filename_prefix):
    """Save content to a text file"""
    os.makedirs("outputs", exist_ok=True)
    content_str = "\n".join(content) if isinstance(content, list) else str(content)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filepath = f"outputs/{filename_prefix}_{timestamp}.txt"
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content_str)
    return filepath

def extract_text_from_file(filepath):
    """Extract text from multiple file formats"""
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"File not found: {filepath}")
    
    ext = os.path.splitext(filepath)[-1].lower()
    try:
        if ext == '.txt':
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
        elif ext == '.pdf':
            with fitz.open(filepath) as doc:
                return "\n".join(page.get_text() for page in doc)
        elif ext == '.docx':
            doc = docx.Document(filepath)
            return "\n".join(para.text for para in doc.paragraphs if para.text)
        elif ext == '.pptx':
            prs = Presentation(filepath)
            return "\n".join(shape.text for slide in prs.slides 
                           for shape in slide.shapes 
                           if hasattr(shape, "text") and shape.text)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
    except Exception as e:
        raise RuntimeError(f"Error processing file: {str(e)}")

def generate_quiz(text):
    """Generate quiz questions and answers"""
    prompt = f"""Create a 5-10 question quiz based on this content. Format exactly as:
    
1. [Question]
   Answer: [Answer]

2. [Question]
   Answer: [Answer]

Content:
{text[:3000]}"""  # Truncate very long texts
    
    response = client.chat.completions.create(
        model="meta-llama/llama-3-70b-instruct",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7,
        max_tokens=2000
    )
    return response.choices[0].message.content

def generate_notes(text):
    """Generate organized study notes"""
    prompt = f"""Create comprehensive study notes with these sections:
    
# Key Concepts
- [Concept 1]
- [Concept 2]

# Important Details
• [Detail 1]
• [Detail 2]

# Examples
✓ [Example 1]
✓ [Example 2]

From this content:
{text[:3000]}"""
    
    response = client.chat.completions.create(
        model="meta-llama/llama-3-70b-instruct",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7
    )
    return response.choices[0].message.content

def generate_mnemonics(words):
    """Create memory aids for word lists"""
    prompt = f"""Create creative mnemonics for these words:
    
{", ".join(words)}

Format each as:
[Word]: [Mnemonic Device] - [Explanation]"""
    
    response = client.chat.completions.create(
        model="meta-llama/llama-3-70b-instruct",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.8
    )
    return response.choices[0].message.content

def generate_story(words):
    """Generate a story using provided words"""
    prompt = f"""Write a creative short story (3-5 paragraphs) using these words:
    
{", ".join(words)}

Story Requirements:
- Use all words naturally
- Include character development
- Have a clear plot"""
    
    response = client.chat.completions.create(
        model="meta-llama/llama-3-70b-instruct",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.9
    )
    return response.choices[0].message.content

def main():
    print("📚 Ultimate Study Tool v2.0")
    print("Supports: PDF, DOCX, PPTX, TXT files\n")
    
    while True:
        print("\n" + "="*40)
        print("Main Menu:")
        print("1. Generate Quiz from File")
        print("2. Create Study Notes")
        print("3. Generate Mnemonics")
        print("4. Create Story")
        print("5. Exit")
        
        choice = input("\nChoose (1-5): ").strip()
        
        if choice == "1":
            filepath = input("\nEnter file path: ").strip('"\' ')
            try:
                text = extract_text_from_file(filepath)
                print("\n⏳ Generating quiz...")
                quiz = generate_quiz(text)
                print("\n✅ Quiz Generated:\n")
                print(quiz)
                
                if input("\n💾 Save quiz? (y/n): ").lower() == 'y':
                    path = save_to_file(quiz, "quiz")
                    print(f"Saved to: {path}")
                    
            except Exception as e:
                print(f"\n❌ Error: {e}")
                
        elif choice == "2":
            print("\nCreate notes from:")
            print("1. File")
            print("2. Direct Text Input")
            sub_choice = input("Choose (1/2): ").strip()
            
            if sub_choice == "1":
                filepath = input("\nEnter file path: ").strip('"\' ')
                try:
                    text = extract_text_from_file(filepath)
                except Exception as e:
                    print(f"\n❌ Error: {e}")
                    continue
            elif sub_choice == "2":
                text = input("\nEnter your content: ").strip()
            else:
                print("\n⚠️ Invalid choice")
                continue
                
            print("\n⏳ Generating notes...")
            notes = generate_notes(text)
            print("\n✅ Notes Generated:\n")
            print(notes)
            
            if input("\n💾 Save notes? (y/n): ").lower() == 'y':
                path = save_to_file(notes, "notes")
                print(f"Saved to: {path}")
                
        elif choice == "3":
            words = [w.strip() for w in input("\nEnter words (comma separated): ").split(",") if w.strip()]
            if not words:
                print("\n⚠️ Please enter at least one word")
                continue
                
            print("\n⏳ Generating mnemonics...")
            mnemonics = generate_mnemonics(words)
            print("\n✅ Mnemonics Generated:\n")
            print(mnemonics)
            
            if input("\n💾 Save mnemonics? (y/n): ").lower() == 'y':
                path = save_to_file(mnemonics, "mnemonics")
                print(f"Saved to: {path}")
                
        elif choice == "4":
            words = [w.strip() for w in input("\nEnter words to include in story: ").split(",") if w.strip()]
            if not words:
                print("\n⚠️ Please enter at least one word")
                continue
                
            print("\n⏳ Generating story...")
            story = generate_story(words)
            print("\n✅ Story Generated:\n")
            print(story)
            
            if input("\n💾 Save story? (y/n): ").lower() == 'y':
                path = save_to_file(story, "story")
                print(f"Saved to: {path}")
                
        elif choice == "5":
            print("\n👋 Goodbye!")
            break
            
        else:
            print("\n⚠️ Please choose 1-5")

if __name__ == "__main__":
    main()
